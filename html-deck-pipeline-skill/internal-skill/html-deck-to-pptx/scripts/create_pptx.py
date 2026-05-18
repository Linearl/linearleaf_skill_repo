"""
将截图文件夹组装为PPTX文件
用法: python create_pptx.py --screenshot-dir <dir> --output <file.pptx> [options]
"""
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu

ASPECT_RATIOS = {
    "16:9": (Inches(13.333), Inches(7.5)),
    "4:3":  (Inches(10),      Inches(7.5)),
    "16:10":(Inches(13.333), Inches(8.333)),
}


def build(screenshot_dir: str, output: str, total: int, aspect: str):
    src = Path(screenshot_dir)
    w, h = ASPECT_RATIOS.get(aspect, ASPECT_RATIOS["16:9"])

    prs = Presentation()
    prs.slide_width = w
    prs.slide_height = h
    layout = prs.slide_layouts[6]  # blank

    added = 0
    for i in range(1, total + 1):
        img = src / f"{i}.jpg"
        if not img.exists():
            print(f"警告：找不到 {img}，跳过")
            continue
        slide = prs.slides.add_slide(layout)
        slide.shapes.add_picture(str(img), Emu(0), Emu(0), w, h)
        added += 1
        print(f"已添加：第 {i} 页")

    prs.save(output)
    print(f"\n完成！{added} 页 → {output}")


def main():
    ap = argparse.ArgumentParser(description="截图文件夹 → PPTX")
    ap.add_argument("--screenshot-dir", required=True, help="截图目录（包含 1.jpg ~ N.jpg）")
    ap.add_argument("--output", required=True, help="输出PPTX路径")
    ap.add_argument("--total", type=int, default=30, help="总页数（默认30）")
    ap.add_argument("--aspect", default="16:9", choices=list(ASPECT_RATIOS.keys()),
                    help="幻灯片比例（默认16:9）")
    args = ap.parse_args()
    build(args.screenshot_dir, args.output, args.total, args.aspect)


if __name__ == "__main__":
    main()
